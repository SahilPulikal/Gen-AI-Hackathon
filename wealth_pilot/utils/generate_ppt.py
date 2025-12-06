from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title, content_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "WealthPilot: The Agentic AI Co-pilot"
    subtitle.text = "Proactive. Intelligent. Client-Centric.\nHackathon Submission"

    # Slide 2: The Challenge
    add_slide("The Challenge: Data Overload", [
        "Wealth managers track thousands of news events daily.",
        "Most client interactions are reactive, not proactive.",
        "Connecting generic news to specific portfolios is manual and slow.",
        "Result: Missed opportunities and lower client trust."
    ])

    # Slide 3: The Solution
    add_slide("The Solution: WealthPilot", [
        "An Agentic AI workspace for proactive wealth management.",
        "Core Philosophy: Don't just show data. Interpret it.",
        "Autonomous Agents work in the background to flag risks.",
        "Empowers managers to be 'Client-First' again."
    ])

    # Slide 4: Key Features
    add_slide("Key Features: Empowering the Manager", [
        "Strategy Crew (NEW): Multi-Agent team for portfolio planning.",
        "Black Swan Simulator (NEW): Market crash impact modeling.",
        "Real-Time Feedback: Transparent AI reasoning UI.",
        "Client 360° Dashboard: Deep risk & holding analysis.",
        "News Impact Agent: Personalized sentiment scanning."
    ])

    # Slide 5: Architecture
    add_slide("Architecture & Tech Stack", [
        "Frontend: Streamlit (Python)",
        "Intelligence: Google Gemini 2.5 Flash",
        "Agent Framework: Custom Multi-Agent Engine (CrewAI-style)",
        "Data Layer: Yahoo Finance (Live) + Synthetic Data"
    ])

    # Slide 6: Business Impact
    add_slide("Business Impact", [
        "Efficiency: Reduce research time by ~40%.",
        "Proactivity: Reach out to clients before they panic.",
        "Personalization: Tailored insights for every single client.",
        "Scalability: Manage more clients with higher quality service."
    ])

    # Slide 7: Conclusion
    add_slide("The Future of Wealth Management", [
        "WealthPilot turns data into action.",
        "Next Steps: CRM Integration & Compliance Checks.",
        "Thank You!"
    ])

    prs.save('WealthPilot_Presentation.pptx')
    print("Presentation saved as WealthPilot_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
